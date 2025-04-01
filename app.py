from pyecharts import options as opts
from pyecharts.charts import Line
from pyecharts.globals import ThemeType
import streamlit as st
import pandas as pd
import os
import math
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
import time
from pptx import Presentation
from pptx.util import Inches
import io

# 定义颜色，使用对比度大且亮度不高的颜色
line_colors = ['#8B0000', '#006400', '#00008B', '#8B8B00', '#8B008B', '#008B8B', '#FF8C00', '#4B0082']

# 缓存 Excel 文件读取，使用 st.cache_resource
@st.cache_resource
def read_excel_file():
    file_path = os.path.join(os.path.dirname(__file__), '沥青数据豆包3.xlsx')
    excel_file = pd.ExcelFile(file_path)
    return excel_file

excel_file = read_excel_file()

# 获取所有表名
sheet_names = excel_file.sheet_names

# Streamlit 布局
st.set_page_config(layout="wide")
col1, col2 = st.columns([1, 3])

# 缓存数据处理函数
@st.cache_data
def process_sheet(sheet_name):
    df = excel_file.parse(sheet_name, header=4)  # 从第 5 行开始加载数据
    headers = excel_file.parse(sheet_name, nrows=6)  # 加载前 6 行作为表头

    # 获取表头信息
    first_row = headers.iloc[0].values
    second_row = headers.iloc[1].values
    fourth_row = headers.iloc[3].values
    fifth_row = headers.iloc[4].values
    sixth_row = headers.iloc[5].values

    date_column = df.columns[0]

    # 处理日期列
    df[date_column] = pd.to_datetime(df[date_column], errors='coerce')

    # 关键修复：处理数值列（强制转换为数字）
    data_columns = df.columns[1:]  # 所有数据列（除日期列）
    for col in data_columns:
        df[col] = pd.to_numeric(df[col], errors='coerce')  # 转换失败设为 NaN

    df = df.dropna(subset=[date_column]).copy()  # 仅删除日期缺失行

    # 创建列名到分类标签的映射
    category_mapping = {}
    for col, label1, label2 in zip(df.columns[1:], first_row[1:], second_row[1:]):
        category_mapping[col] = (label1, label2)

    category_labels1 = set([label[0] for label in category_mapping.values()])

    return df, date_column, category_mapping, fourth_row, fifth_row, sixth_row

# 新增辅助函数
def calculate_yaxis_limits(data, padding_ratio=0.05):
    """计算带扩展范围的整数坐标轴"""
    valid_data = data.dropna()
    if valid_data.empty:
        return 0, 1, 1

    data_min = valid_data.min()
    data_max = valid_data.max()
    data_range = data_max - data_min

    # 计算扩展范围
    padding = padding_ratio * data_range
    data_min -= padding
    data_max += padding

    # 取整到合适的整数范围
    min_round = math.floor(data_min)
    max_round = math.ceil(data_max)

    # 计算合理的间隔
    interval = max(1, (max_round - min_round) // 5)

    return min_round, max_round, interval

# 缓存时间序列图函数
@st.cache_data
def create_time_series_chart(df, date_column, selected_column):
    single_df = df[[date_column, selected_column]].dropna()
    single_df = single_df.sort_values(by=date_column)

    # 计算坐标轴范围
    y_min, y_max, interval = calculate_yaxis_limits(single_df[selected_column])

    line = (
        Line(init_opts=opts.InitOpts(theme=ThemeType.LIGHT, width="1000px", height="800px", bg_color="white"))  # 明确设置背景颜色
        .add_xaxis(single_df[date_column].dt.strftime('%Y-%m-%d').tolist())
        .add_yaxis(selected_column, single_df[selected_column].tolist(), is_smooth=True,
                   label_opts=opts.LabelOpts(is_show=False))
        .set_global_opts(
            title_opts=opts.TitleOpts(title=f"{selected_column} 时间序列图"),
            toolbox_opts=opts.ToolboxOpts(is_show=True),
            xaxis_opts=opts.AxisOpts(name="日期"),
            yaxis_opts=opts.AxisOpts(
                name=selected_column,
                min_=y_min,
                max_=y_max,
                interval=interval
            ),
            datazoom_opts=[opts.DataZoomOpts(type_="slider", xaxis_index=0, range_start=0, range_end=100)]
        )
    )
    return line

# 缓存季节性图表函数
@st.cache_data
def create_seasonal_chart(df, date_column, selected_column, fourth_row, fifth_row, selected_year_range):
    single_df = df[[date_column, selected_column]].dropna()
    single_df['年份'] = single_df[date_column].dt.year.astype(int)
    years = sorted(single_df['年份'].unique(), reverse=True)

    if selected_year_range == "5年":
        years = years[:5]
    elif selected_year_range == "8年":
        years = years[:8]

    if '日' in fifth_row[list(fourth_row).index(selected_column)]:
        single_df['日序'] = single_df[date_column].dt.dayofyear
        x_axis_name = '日序（1 - 366）'
        x_axis_values = list(range(1, 367))
    elif '周' in fifth_row[list(fourth_row).index(selected_column)]:
        single_df['周序'] = single_df[date_column].dt.isocalendar().week
        x_axis_name = '周序（1 - 53）'
        x_axis_values = list(range(1, 54))
    else:  # 月频率
        single_df['月序'] = single_df[date_column].dt.month
        x_axis_name = '月份（1 - 12）'
        x_axis_values = list(range(1, 13))

    all_y_values = []
    for year in years:
        if '日' in fifth_row[list(fourth_row).index(selected_column)]:
            year_data = single_df[(single_df['年份'] == year)].groupby('日序')[selected_column].mean()
        elif '周' in fifth_row[list(fourth_row).index(selected_column)]:
            year_data = single_df[(single_df['年份'] == year)].groupby('周序')[selected_column].mean()
        else:
            year_data = single_df[(single_df['年份'] == year)].groupby('月序')[selected_column].mean()
        y_values = [year_data.get(x, None) for x in x_axis_values]
        all_y_values.extend([y for y in y_values if y is not None])

    # 计算坐标轴范围
    y_min, y_max, interval = calculate_yaxis_limits(pd.Series(all_y_values))

    line = (
        Line(init_opts=opts.InitOpts(theme=ThemeType.LIGHT, width="1000px", height="800px", bg_color="white"))  # 明确设置背景颜色
        .set_global_opts(
            title_opts=opts.TitleOpts(title=f"{selected_column} 季节性图表"),
            toolbox_opts=opts.ToolboxOpts(is_show=True),
            xaxis_opts=opts.AxisOpts(name=x_axis_name),
            yaxis_opts=opts.AxisOpts(
                name=selected_column,
                min_=y_min,
                max_=y_max,
                interval=interval
            ),
            legend_opts=opts.LegendOpts(is_show=True, type_="scroll", pos_bottom="1%", pos_left="center")
        )
    )

    custom_colors = ['#FF0000', '#000000', '#0000FF', '#00FF00', '#800080', '#FFA500']
    for i, year in enumerate(years):
        if '日' in fifth_row[list(fourth_row).index(selected_column)]:
            year_data = single_df[(single_df['年份'] == year)].groupby('日序')[selected_column].mean()
        elif '周' in fifth_row[list(fourth_row).index(selected_column)]:
            year_data = single_df[(single_df['年份'] == year)].groupby('周序')[selected_column].mean()
        else:
            year_data = single_df[(single_df['年份'] == year)].groupby('月序')[selected_column].mean()
        y_values = [year_data.get(x, None) for x in x_axis_values]
        line.add_xaxis(x_axis_values)

        if i < len(custom_colors):
            color = custom_colors[i]
        else:
            color = line_colors[(i - len(custom_colors)) % len(line_colors)]

        line.add_yaxis(
            str(year),
            y_values,
            is_smooth=True,
            label_opts=opts.LabelOpts(is_show=False),
            linestyle_opts=opts.LineStyleOpts(color=color, width=3)
        )
    return line

# HTML 转 PNG 函数
def html_to_png(html_content, output_path):
    temp_html_path = None
    driver = None
    try:
        # 指定 ChromeDriver 版本，需要根据实际 Chrome 版本修改
        service = Service(ChromeDriverManager(version="114.0.5735.90").install())
        driver = webdriver.Chrome(service=service)
        # 将 HTML 内容保存为临时文件
        temp_html_path = os.path.join(os.getcwd(), "temp_chart.html")
        with open(temp_html_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        # 打开临时 HTML 文件
        driver.get(f"file:///{temp_html_path}")

        # 等待页面加载完成
        time.sleep(5)

        # 保存为 PNG
        driver.save_screenshot(output_path)
        print(f"成功生成 PNG 文件: {output_path}")
    except Exception as e:
        print(f"Error converting HTML to PNG: {e}")
    finally:
        if driver:
            # 关闭浏览器
            driver.quit()
        if temp_html_path and os.path.exists(temp_html_path):
            # 删除临时 HTML 文件
            os.remove(temp_html_path)

with col1:
    # 选择 Sheet 名
    selected_sheet = st.selectbox("选择 Sheet 名", sheet_names)
    df, date_column, category_mapping, fourth_row, fifth_row, sixth_row = process_sheet(selected_sheet)

    # 选择第一行数据标签
    category_labels1 = set([label[0] for label in category_mapping.values()])
    selected_category1 = st.selectbox(f"选择 {selected_sheet} 的第一行数据标签", category_labels1)

    # 根据第一行标签筛选第二行细分标签
    available_labels2 = [label[1] for col, label in category_mapping.items() if label[0] == selected_category1]
    selected_category2 = st.selectbox(f"选择 {selected_sheet} 的第二行细分标签", set(available_labels2))

    # 根据前两个选择筛选指标名称，使用 st.multiselect 支持多选
    available_columns = [col for col, label in category_mapping.items() if
                         label[0] == selected_category1 and label[1] == selected_category2]
    selected_columns = st.multiselect(f"选择 {selected_sheet} 的指标名称", available_columns)

    # 选择图表类型
    chart_type = st.selectbox(f"选择 {selected_sheet} 的图表类型", ["时间序列图", "季节性图表"])

    if chart_type == "季节性图表":
        # 新增：选择年份范围
        year_range_options = ["5年", "8年", "全部"]
        selected_year_range = st.selectbox("选择展示的年份范围", year_range_options, index=0)

    # 导出 PPT 按钮
    if st.button("导出选中图表到 PPT"):
        prs = Presentation()
        for selected_column in selected_columns:
            if chart_type == "季节性图表":
                chart = create_seasonal_chart(df, date_column, selected_column, fourth_row, fifth_row, selected_year_range)
            else:
                chart = create_time_series_chart(df, date_column, selected_column)

            html_content = chart.render_embed()
            output_path = os.path.join(os.getcwd(), f"{selected_column}_{chart_type}.png")
            html_to_png(html_content, output_path)

            if os.path.exists(output_path):
                # 创建新的幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                body_shape = slide.shapes.placeholders[1]

                # 设置幻灯片标题
                title.text = f"{selected_column} {chart_type}"

                # 添加图表图片到幻灯片
                left = Inches(1)
                top = Inches(1.5)
                try:
                    pic = slide.shapes.add_picture(output_path, left, top, width=Inches(8), height=Inches(6))
                except Exception as e:
                    print(f"添加图片到幻灯片时出错: {e}")

                # 显示数据描述
                description = sixth_row[list(fourth_row).index(selected_column)]
                tf = body_shape.text_frame
                tf.text = f"数据描述：{description}"
            else:
                print(f"未找到 PNG 文件: {output_path}")

        # 保存 PPT 文件
        buffer = io.BytesIO()
        try:
            prs.save(buffer)
            buffer.seek(0)
            # 提供下载链接
            st.download_button(
                label="下载 PPT",
                data=buffer,
                file_name="charts.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            print(f"保存 PPT 到缓冲区时出错: {e}")

with col2:
    for selected_column in selected_columns:
        if chart_type == "季节性图表":
            chart = create_seasonal_chart(df, date_column, selected_column, fourth_row, fifth_row, selected_year_range)
        else:
            chart = create_time_series_chart(df, date_column, selected_column)
        html_content = chart.render_embed()
        st.components.v1.html(html_content, height=800)

        # 显示数据描述
        description = sixth_row[list(fourth_row).index(selected_column)]
        st.markdown(f"<small>数据描述：{description}</small>", unsafe_allow_html=True)

        # 转换为 PNG
        output_path = os.path.join(os.getcwd(), f"{selected_column}_{chart_type}.png")
        html_to_png(html_content, output_path)
        if os.path.exists(output_path):
            st.success(f"已将 {selected_column} 的 {chart_type} 转换为 PNG: {output_path}")
        else:
            st.error(f"未能将 {selected_column} 的 {chart_type} 转换为 PNG")
